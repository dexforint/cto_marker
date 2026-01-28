"""
Провайдер для обработки презентаций PowerPoint (PPTX).

Модуль предоставляет PowerPointProvider, который конвертирует презентации
в формат PDF для последующей обработки единообразным способом.
Процесс включает извлечение содержимого слайдов (текст, таблицы, изображения)
и их конвертацию в HTML структуру с последующим преобразованием в PDF.

Автор: Marker Team
"""

import base64
import os
import tempfile
import traceback

from marker.logger import get_logger
from marker.providers.pdf import PdfProvider

# Инициализируем логгер для записи ошибок и отладочной информации
logger = get_logger()

# CSS стили для отображения презентаций в альбомной ориентации
css = """
@page {
    size: A4 landscape;  # Альбомная ориентация страницы A4
    margin: 1.5cm;  # Отступы по 1.5 см со всех сторон
}

table {
    width: 100%;  # Ширина таблицы 100%
    border-collapse: collapse;  # Слияние границ ячеек
    break-inside: auto;  # Автоматический перенос таблиц
    font-size: 10pt;  # Размер шрифта 10 пунктов
}

tr {
    break-inside: avoid;  # Избегать переноса строк таблицы
    page-break-inside: avoid;  # Избегать разрыва страницы внутри строки
}

td {
    border: 0.75pt solid #000;  # Границы ячеек: 0.75pt, сплошные, черные
    padding: 6pt;  # Внутренние отступы ячеек
}

img {
    max-width: 100%;  # Максимальная ширина изображения 100%
    height: auto;  # Автоматическая высота для сохранения пропорций
    object-fit: contain;  # Сохранение пропорций изображения
}
"""


class PowerPointProvider(PdfProvider):
    """
    Провайдер для обработки презентаций PowerPoint (PPTX).
    
    Наследуется от PdfProvider и обеспечивает конвертацию презентаций
    в промежуточный формат PDF для последующей обработки.
    
    Процесс обработки:
    1. Создание временного PDF файла
    2. Извлечение содержимого слайдов (текст, таблицы, изображения)
    3. Конвертация в HTML структуру с сохранением форматирования
    4. Конвертация HTML в PDF
    5. Инициализация родительского PdfProvider с временным PDF
    
    Атрибуты:
        include_slide_number (bool): Включать ли номера слайдов в выходной документ
    """
    # Флаг для включения номеров слайдов в HTML вывод
    include_slide_number: bool = False

    def __init__(self, filepath: str, config=None):
        """
        Инициализация провайдера презентаций.
        
        Args:
            filepath (str): Путь к PPTX файлу
            config: Конфигурация провайдера (опционально)
        """
        # Создаем временный PDF файл для промежуточного хранения
        temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        self.temp_pdf_path = temp_pdf.name
        temp_pdf.close()

        # Конвертируем PPTX в PDF
        try:
            self.convert_pptx_to_pdf(filepath)
        except Exception as e:
            # Выводим подробную информацию об ошибке и выбрасываем исключение
            print(traceback.format_exc())
            raise ValueError(f"Error converting PPTX to PDF: {e}")

        # Инициализируем родительский PdfProvider с временным PDF файлом
        super().__init__(self.temp_pdf_path, config)

    def __del__(self):
        """
        Деструктор для очистки временных файлов.
        
        Удаляет временный PDF файл при уничтожении объекта.
        """
        if os.path.exists(self.temp_pdf_path):
            os.remove(self.temp_pdf_path)

    def convert_pptx_to_pdf(self, filepath):
        """
        Конвертирует презентацию PPTX в PDF формат.
        
        Процесс включает:
        1. Чтение PPTX файла с помощью python-pptx
        2. Обработку каждого слайда и его элементов
        3. Конвертацию различных типов элементов в HTML
        4. Объединение всех HTML частей
        5. Конвертацию HTML в PDF
        
        Args:
            filepath (str): Путь к исходному PPTX файлу
        """
        from weasyprint import CSS, HTML
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Загружаем презентацию
        pptx = Presentation(filepath)

        # Список для хранения частей HTML
        html_parts = []

        # Обрабатываем каждый слайд
        for slide_index, slide in enumerate(pptx.slides):
            # Начинаем новую секцию для слайда
            html_parts.append("<section>")
            # Добавляем заголовок слайда, если включена нумерация
            if self.include_slide_number:
                html_parts.append(f"<h2>Slide {slide_index + 1}</h2>")

            # Обрабатываем фигуры на слайде
            for shape in slide.shapes:
                # Если фигура является группой, обрабатываем рекурсивно
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    html_parts.append(self._handle_group(shape))
                    continue

                # Если фигура содержит таблицу
                if shape.has_table:
                    html_parts.append(self._handle_table(shape))
                    continue

                # Если фигура является изображением
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    html_parts.append(self._handle_image(shape))
                    continue

                # Если фигура содержит текст
                if hasattr(shape, "text") and shape.text is not None:
                    if shape.has_text_frame:
                        # Различаем плейсхолдеры (заголовок, подзаголовок и т.д.)
                        html_parts.append(self._handle_text(shape))
                    else:
                        # Обычный текст без текстового фрейма
                        html_parts.append(f"<p>{self._escape_html(shape.text)}</p>")

            # Завершаем секцию слайда
            html_parts.append("</section>")

        # Объединяем все части HTML в одну строку
        html = "\n".join(html_parts)

        # Конвертируем HTML в PDF с применением стилей
        HTML(string=html).write_pdf(
            self.temp_pdf_path, stylesheets=[CSS(string=css), self.get_font_css()]
        )

    def _handle_group(self, group_shape) -> str:
        """
        Рекурсивно обрабатывает фигуры в группе.
        
        Функция проходит по всем фигурам в группе и обрабатывает каждую
        в зависимости от ее типа (группа, таблица, изображение, текст).
        
        Args:
            group_shape: Объект группы фигур из python-pptx
            
        Returns:
            str: HTML строка для всей группы фигур
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        group_parts = []
        for shape in group_shape.shapes:
            # Если фигура является группой, обрабатываем рекурсивно
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_parts.append(self._handle_group(shape))
                continue

            # Если фигура содержит таблицу
            if shape.has_table:
                group_parts.append(self._handle_table(shape))
                continue

            # Если фигура является изображением
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                group_parts.append(self._handle_image(shape))
                continue

            # Если фигура содержит текст
            if hasattr(shape, "text"):
                if shape.has_text_frame:
                    group_parts.append(self._handle_text(shape))
                else:
                    group_parts.append(f"<p>{self._escape_html(shape.text)}</p>")

        return "".join(group_parts)

    def _handle_text(self, shape) -> str:
        """
        Обрабатывает текст фигуры, включая определение списков и плейсхолдеров.
        
        Функция определяет тип текстового элемента (заголовок, подзаголовок,
        список) и конвертирует его в соответствующие HTML теги.
        
        Args:
            shape: Объект фигуры с текстом из python-pptx
            
        Returns:
            str: HTML строка для текстового блока
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        # Определяем HTML тег на основе типа плейсхолдера
        label_html_tag = "p"
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                label_html_tag = "h3"
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                label_html_tag = "h4"

        # Следим за тем, находимся ли мы в списке <ul> или <ol>
        html_parts = []
        list_open = False
        list_type = None  # "ul" или "ol"

        # Обрабатываем каждый параграф в текстовом фрейме
        for paragraph in shape.text_frame.paragraphs:
            p_el = paragraph._element
            # Проверяем маркеры списков
            bullet_char = p_el.find(".//a:buChar", namespaces=p_el.nsmap)
            bullet_num = p_el.find(".//a:buAutoNum", namespaces=p_el.nsmap)

            # Определяем тип списка
            is_bullet = (bullet_char is not None) or (paragraph.level > 0)
            is_numbered = bullet_num is not None

            # Если параграф является элементом списка
            if is_bullet or is_numbered:
                # Определяем тип текущего списка
                current_list_type = "ol" if is_numbered else "ul"
                if not list_open:
                    # Начинаем новый список
                    list_open = True
                    list_type = current_list_type
                    html_parts.append(f"<{list_type}>")

                elif list_open and list_type != current_list_type:
                    # Закрываем старый список и начинаем новый
                    html_parts.append(f"</{list_type}>")
                    list_type = current_list_type
                    html_parts.append(f"<{list_type}>")

                # Собираем текст элемента списка из всех runs
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text:
                    html_parts.append(f"<li>{self._escape_html(p_text)}</li>")

            else:
                # Если мы были в списке, закрываем его
                if list_open:
                    html_parts.append(f"</{list_type}>")
                    list_open = False
                    list_type = None

                # Обычный параграф
                # Собираем текст параграфа из всех runs
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text:
                    # Используем соответствующий HTML тег
                    html_parts.append(
                        f"<{label_html_tag}>{self._escape_html(p_text)}</{label_html_tag}>"
                    )

        # Если текстовый фрейм закончился, а список все еще открыт, закрываем его
        if list_open:
            html_parts.append(f"</{list_type}>")

        return "".join(html_parts)

    def _handle_image(self, shape) -> str:
        """
        Встраивает изображение как base64 <img> в HTML.
        
        Извлекает изображение из фигуры презентации и конвертирует его
        в base64 строку для встраивания в HTML.
        
        Args:
            shape: Объект фигуры с изображением из python-pptx
            
        Returns:
            str: HTML тег img с base64 изображением или пустая строка при ошибке
        """
        image = shape.image
        image_bytes = image.blob

        try:
            # Кодируем изображение в base64
            img_str = base64.b64encode(image_bytes).decode("utf-8")
            return f"<img src='data:{image.content_type};base64,{img_str}' />"
        except Exception as e:
            # Логируем предупреждение при ошибке загрузки изображения
            logger.warning(f"Warning: image cannot be loaded by Pillow: {e}")
            return ""

    def _handle_table(self, shape) -> str:
        """
        Отображает таблицу фигуры как HTML <table>.
        
        Конвертирует таблицу из презентации в HTML таблицу с границами.
        
        Args:
            shape: Объект фигуры с таблицей из python-pptx
            
        Returns:
            str: HTML строка с таблицей
        """
        table_html = []
        table_html.append("<table border='1'>")

        # Обрабатываем каждую строку таблицы
        for row in shape.table.rows:
            row_html = ["<tr>"]
            # Обрабатываем каждую ячейку в строке
            for cell in row.cells:
                row_html.append(f"<td>{self._escape_html(cell.text)}</td>")
            row_html.append("</tr>")
            table_html.append("".join(row_html))

        table_html.append("</table>")
        return "".join(table_html)

    def _escape_html(self, text: str) -> str:
        """
        Минимальное экранирование HTML специальных символов.
        
        Заменяет HTML метасимволы на их эквиваленты для безопасного отображения.
        
        Args:
            text (str): Исходный текст для экранирования
            
        Returns:
            str: Экранированный текст
        """
        return (
            text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
        )
